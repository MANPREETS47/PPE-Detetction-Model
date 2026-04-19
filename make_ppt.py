from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_ppt():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # 0 is Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Real-Time PPE Detection System"
    subtitle.text = "Ensuring Workplace Safety with Custom YOLO Object Detection\n\nProject Presentation"

    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Problem: Failing to wear required PPE in industrial environments can lead to severe injuries or fatalities."
    p = tf.add_paragraph()
    p.text = "Current Limitations: Manual monitoring is error-prone, labor-intensive, and impossible to scale 24/7."
    p = tf.add_paragraph()
    p.text = "The Goal: Automate the monitoring process using Artificial Intelligence to ensure continuous safety compliance."

    # Slide 3: Proposed Solution
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Proposed Solution"
    tf = slide.placeholders[1].text_frame
    tf.text = "A real-time Computer Vision application using a custom YOLO model to detect if workers are wearing required safety gear."
    p = tf.add_paragraph()
    p.text = "Analyzes live video feeds (webcam or CCTV)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automatically detects missing gear (Helmet, Vest, Gloves, Goggles, Boots)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Auto-saves violation snapshots with timestamps for auditing."
    p.level = 1

    # Slide 4: Key Features
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Features"
    tf = slide.placeholders[1].text_frame
    tf.text = "Real-Time Detection: Rapid inference frame-by-frame using YOLOv8."
    p = tf.add_paragraph()
    p.text = "Automated Violation Logging: Captures and saves image evidence directly to a violations/ folder."
    p = tf.add_paragraph()
    p.text = "Smart Cooldown System: Implements a 5-second cooldown to prevent disk-spamming when the same violation happens continuously."
    p = tf.add_paragraph()
    p.text = "Flexible Inputs: Works with live webcams or pre-recorded video files."

    # Slide 5: Tech Stack & Tools
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Tech Stack & Tools"
    tf = slide.placeholders[1].text_frame
    tf.text = "Programming Language: Python 3.8+"
    p = tf.add_paragraph()
    p.text = "Deep Learning Framework: Ultralytics YOLOv8 (for object detection)"
    p = tf.add_paragraph()
    p.text = "Computer Vision Library: OpenCV (opencv-python) for grabbing frames and rendering bounding boxes."
    p = tf.add_paragraph()
    p.text = "Array Processing: NumPy"

    # Slide 6: How It Works - The Methodology
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "How It Works - Methodology"
    tf = slide.placeholders[1].text_frame
    tf.text = "Negative Classes Logic: The model explicitly identifies when PPE should be there but isn't (e.g., no_boots, no_helmet) to reduce false positives."
    p = tf.add_paragraph()
    p.text = "If body parts are out of frame, the model won't trigger a false alarm."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The Vest Logic: Because there isn't a 'no_vest' class, the application contains built-in logic:"
    p = tf.add_paragraph()
    p.text = "If 'Person' is detected BUT 'Vest' is NOT detected -> Trigger Vest Violation."
    p.level = 1

    # Slide 7: System Workflow / Architecture
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Workflow / Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Input: Read frame from Webcam/Video."
    p = tf.add_paragraph()
    p.text = "2. Inference: YOLO model (best.pt) analyzes the frame."
    p = tf.add_paragraph()
    p.text = "3. Logic Gate: Filter detected classes. Check EXPLICIT_VIOLATION_MAP and Vest conditions."
    p = tf.add_paragraph()
    p.text = "4. Action: If missing gear -> Draw red alerts on the frame -> Check Cooldown -> Save Image."
    p = tf.add_paragraph()
    p.text = "5. Output: Display annotated frame on screen with running Total Violations count."

    # Slide 8: Future Improvements
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Future Improvements"
    tf = slide.placeholders[1].text_frame
    tf.text = "Integration with a Cloud Dashboard (AWS/Azure) to send email/SMS alerts to safety managers."
    p = tf.add_paragraph()
    p.text = "Multi-camera support for large-scale warehouse deployment."
    p = tf.add_paragraph()
    p.text = "Expanding the dataset to improve model accuracy in low-light conditions."

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion & Q&A"
    tf = slide.placeholders[1].text_frame
    tf.text = "The project successfully demonstrates how Edge AI can be deployed simply and effectively to prevent workplace hazards."
    p = tf.add_paragraph()
    p.text = "Thank You!"
    p = tf.add_paragraph()
    p.text = "Any Questions?"

    # Save the presentation
    ppt_path = "PPE_Detection_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation successfully saved as {ppt_path}")

if __name__ == "__main__":
    create_ppt()
